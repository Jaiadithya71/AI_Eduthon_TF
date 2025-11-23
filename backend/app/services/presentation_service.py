import os
import io
import time
import uuid
from datetime import datetime
from typing import List

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.models.schemas import (
    PresentationRequest,
    PresentationResponse,
    SlideContent,
    SlideType,
    Language,
    PresentationStyle,
)
from app.services.llm_service import llm_service
from app.services.image_service import image_service

# Base directory for saved PPTX files
BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
PRESENTATIONS_DIR = os.path.join(BASE_DIR, "generated_presentations")
os.makedirs(PRESENTATIONS_DIR, exist_ok=True)


class PresentationService:
    """
    Orchestrates:
    - LLM slide generation
    - Hybrid image selection (AI diagrams + Pexels)
    - PowerPoint building with themes and safe layout
    """

    # ------------------------------------------------------------------
    # PUBLIC API
    # ------------------------------------------------------------------
    async def generate_presentation(
        self, request: PresentationRequest
    ) -> PresentationResponse:
        start_time = time.perf_counter()

        # Basic extra safety
        ok, msg = self.validate_request(request)
        if not ok:
            raise ValueError(msg)

        # 1) Generate slide content using the LLM
        slides: List[SlideContent] = llm_service.generate_slide_content(request)

        # 2) Attach hybrid images (and avoid duplicates)
        self._attach_images_to_slides(request, slides)

        # 3) Prepare file path
        presentation_id = f"pres_{uuid.uuid4().hex[:12]}"
        filename = f"eduslide_ai_{presentation_id}.pptx"
        file_path = os.path.join(PRESENTATIONS_DIR, filename)

        # 4) Build PPTX file
        self._build_pptx(file_path, request, slides)

        # 5) Metadata
        generation_time = time.perf_counter() - start_time
        created_at = datetime.utcnow().isoformat()

        return PresentationResponse(
            presentation_id=presentation_id,
            metadata=request,
            slides=slides,
            total_slides=len(slides),
            created_at=created_at,
            generation_time=generation_time,
        )

    # ------------------------------------------------------------------
    # VALIDATION
    # ------------------------------------------------------------------
    def validate_request(self, request: PresentationRequest) -> tuple[bool, str]:
        topic = (request.topic or "").strip()
        if len(topic) < 5:
            return False, "Topic must be at least 5 characters long"

        if request.num_slides < 3 or request.num_slides > 15:
            return False, "Number of slides must be between 3 and 15"

        return True, ""

    # ------------------------------------------------------------------
    # THEME SELECTION  (no external .pptx files needed)
    # ------------------------------------------------------------------
    def _get_theme_for_request(self, request: PresentationRequest) -> dict:
        """
        Return a simple theme config (colors + sizing tweaks)
        based on the selected PresentationStyle.
        """

        # Default "clean light" theme
        theme = {
            "bg": RGBColor(245, 246, 250),
            "title_color": RGBColor(20, 20, 20),
            "body_color": RGBColor(30, 30, 30),
            "accent": RGBColor(52, 120, 246),      # blue
            "subtitle_color": RGBColor(90, 90, 90),
        }

        style = request.presentation_style

        if style == PresentationStyle.ACADEMIC:
            theme = {
                "bg": RGBColor(250, 252, 255),
                "title_color": RGBColor(15, 40, 80),
                "body_color": RGBColor(35, 35, 45),
                "accent": RGBColor(52, 84, 209),    # academic blue
                "subtitle_color": RGBColor(70, 80, 110),
            }
        elif style == PresentationStyle.STORYTELLING:
            theme = {
                "bg": RGBColor(255, 251, 245),
                "title_color": RGBColor(80, 40, 20),
                "body_color": RGBColor(55, 45, 40),
                "accent": RGBColor(230, 126, 34),   # warm orange
                "subtitle_color": RGBColor(120, 90, 70),
            }
        elif style == PresentationStyle.VISUAL:
            theme = {
                "bg": RGBColor(245, 248, 255),
                "title_color": RGBColor(25, 25, 35),
                "body_color": RGBColor(40, 40, 50),
                "accent": RGBColor(46, 204, 113),   # green accent
                "subtitle_color": RGBColor(90, 100, 120),
            }
        elif style == PresentationStyle.TECHNICAL:
            theme = {
                "bg": RGBColor(20, 24, 31),
                "title_color": RGBColor(236, 240, 241),
                "body_color": RGBColor(221, 230, 234),
                "accent": RGBColor(52, 152, 219),   # tech blue
                "subtitle_color": RGBColor(171, 178, 185),
            }

        return theme

    # ------------------------------------------------------------------
    # IMAGE SELECTION
    # ------------------------------------------------------------------
    def _attach_images_to_slides(
        self,
        request: PresentationRequest,
        slides: List[SlideContent],
    ) -> None:
        """
        For each slide, ask the hybrid image engine for the best image URL.
        Also keeps track of used URLs to prevent repetition.
        """
        topic = request.topic
        used_urls: set[str] = set()

        for idx, slide in enumerate(slides):
            url = image_service.get_hybrid_image_for_slide(
                topic=topic,
                slide=slide,
                slide_index=idx,
                language=request.language,
                presentation_style=request.presentation_style,
                used_urls=used_urls,
            )
            slide.image_url = url
            if url:
                used_urls.add(url)

    # ------------------------------------------------------------------
    # PPTX BUILDING (themes + layout + bullet safety)
    # ------------------------------------------------------------------
    def _build_pptx(
        self,
        file_path: str,
        request: PresentationRequest,
        slides: List[SlideContent],
    ) -> None:
        """
        Builds the .pptx file with:
        - 16:9 widescreen size
        - Themed background + colors
        - Title at the top
        - Bullet text on the left (overflow-safe)
        - Image on the right (scaled + centered)
        """

        prs = Presentation()

        # Force 16:9 widescreen so width isn't tiny
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        theme = self._get_theme_for_request(request)
        is_dark = request.presentation_style == PresentationStyle.TECHNICAL
        is_bilingual = request.language == Language.BILINGUAL

        # Text sizing rules
        if is_bilingual:
            max_bullets = 4
            max_chars_per_bullet = 160
            bullet_font_size = Pt(16)
        else:
            max_bullets = 5
            max_chars_per_bullet = 200
            bullet_font_size = Pt(18)

        for slide_index, slide_data in enumerate(slides):
            layout = prs.slide_layouts[6]  # blank slide
            slide = prs.slides.add_slide(layout)

            # -------------------------------
            # BACKGROUND COLOR (theme)
            # -------------------------------
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = theme["bg"]

            # -------------------------------
            # TITLE
            # -------------------------------
            title_box = slide.shapes.add_textbox(
                Inches(0.6),
                Inches(0.3),
                Inches(11.8),
                Inches(1.0),
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_tf.clear()

            title_para = title_tf.paragraphs[0]
            title_para.text = slide_data.title or request.topic
            title_para.font.size = Pt(34)
            title_para.font.bold = True
            title_para.font.color.rgb = theme["title_color"]

            # Optional subtle accent bar under title
            accent_left = Inches(0.6)
            accent_top = Inches(1.25)
            accent_width = Inches(2.5)
            accent_height = Inches(0.08)
            shape = slide.shapes.add_shape(
                autoshape_type_id=1,  # rectangle
                left=accent_left,
                top=accent_top,
                width=accent_width,
                height=accent_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["accent"]
            shape.line.fill.background()

            # -------------------------------
            # PREP BULLETS (truncate to avoid overflow)
            # -------------------------------
            bullets_raw = slide_data.content or []
            if isinstance(bullets_raw, str):
                bullets_raw = [bullets_raw]

            cleaned: List[str] = []
            for b in bullets_raw:
                if not b:
                    continue
                text = str(b).strip()
                if not text:
                    continue
                if len(text) > max_chars_per_bullet:
                    cut = text[:max_chars_per_bullet]
                    if " " in cut:
                        cut = cut.rsplit(" ", 1)[0]
                    text = cut + "..."
                cleaned.append(text)

            if len(cleaned) > max_bullets:
                head = cleaned[: max_bullets - 1]
                tail = cleaned[max_bullets - 1 :]
                merged_tail = "; ".join(tail)
                if len(merged_tail) > max_chars_per_bullet:
                    merged_cut = merged_tail[:max_chars_per_bullet]
                    if " " in merged_cut:
                        merged_cut = merged_cut.rsplit(" ", 1)[0]
                    merged_tail = merged_cut + "..."
                head.append("Further details: " + merged_tail)
                cleaned = head

            bullets = cleaned
            if not bullets and slide_data.type != SlideType.TITLE:
                bullets = [f"Key ideas about {slide_data.title or request.topic}."]

            # -------------------------------
            # TEXT BLOCK (LEFT) – dynamic vertical position
            # -------------------------------
            bullet_count = len(bullets)

            if slide_data.type == SlideType.TITLE and not bullets_raw:
                # Pure title slide (no bullets)
                text_top = Inches(2.3)
                text_height = Inches(3.5)
            else:
                if bullet_count <= 3:
                    # Short slide: center text more vertically
                    text_top = Inches(2.1)
                    text_height = Inches(3.8)
                else:
                    # Normal / long content
                    text_top = Inches(1.5)
                    text_height = Inches(5.1)

            content_box = slide.shapes.add_textbox(
                Inches(0.8),
                text_top,
                Inches(6.1),
                text_height,
            )
            tf = content_box.text_frame
            tf.clear()
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = f"• {bullet}"
                p.level = 0
                p.font.size = bullet_font_size
                p.font.bold = False
                p.font.color.rgb = theme["body_color"]
                p.line_spacing = 1.15
                p.space_after = Pt(4)
                p.space_before = Pt(1)

            # -------------------------------
            # IMAGE BLOCK (RIGHT)
            # -------------------------------
            if slide_data.image_url:
                img_stream = self._download_image(slide_data.image_url)

                if img_stream:
                    try:
                        img_left = Inches(7.2)
                        img_top = Inches(1.7)
                        max_width = Inches(5.6)
                        max_height = Inches(5.0)

                        pic = slide.shapes.add_picture(img_stream, img_left, img_top)

                        # Scale to fit width
                        if pic.width > max_width:
                            scale = max_width / pic.width
                            pic.width = int(pic.width * scale)
                            pic.height = int(pic.height * scale)

                        # Scale to fit height
                        if pic.height > max_height:
                            scale = max_height / pic.height
                            pic.width = int(pic.width * scale)
                            pic.height = int(pic.height * scale)

                        # Center vertically in reserved block
                        pic.top = int(img_top + (max_height - pic.height) / 2)

                    except Exception as e:
                        print(f"[PPTX] Error placing image: {e}")

        prs.save(file_path)

    # ------------------------------------------------------------------
    # IMAGE DOWNLOAD
    # ------------------------------------------------------------------
    def _download_image(self, url: str) -> io.BytesIO | None:
        try:
            resp = requests.get(url, timeout=12)
            resp.raise_for_status()
            return io.BytesIO(resp.content)
        except Exception as e:
            print(f"[ImageDownload] Error for '{url}': {e}")
            return None


# Singleton instance
presentation_service = PresentationService()
